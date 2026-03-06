#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Created on Mon Mar  3 15:05:47 2025

@author: mikewoodward
"""

from pptx import Presentation
from pptx.enum.shapes import (PP_PLACEHOLDER, 
                              MSO_SHAPE_TYPE)
    
presentation = Presentation(
  "Presentation1.pptx")

# Extract the properties
print("""Slide properties""")
print("""================""")
print(
  "Title: "
  f"{presentation.core_properties.title}")
print(
  "Author: "
  f"{presentation.core_properties.author}")
print(
  "Created: "
  f"{presentation.core_properties.created}")
print(
  "Modified: "
  f"{presentation.core_properties.modified}")

# Find and extract the footer
print("""Slide footers""")
print("""=============""")
for slide in presentation.slides:
  for shape in slide.shapes:
    if (shape.is_placeholder and 
        shape.placeholder_format.type 
        == PP_PLACEHOLDER.FOOTER):
      footer_text = shape.text_frame.text
      slide_no = (presentation
                  .slides
                  .index(slide) + 1)
      print(f"Slide: {slide_no}")
      print(f"Footer: {footer_text}")

# Find and extract all images
print("""Slide images""")
print("""============""")
for slide in presentation.slides:
  for shape in slide.shapes:
    if (shape.shape_type 
        == MSO_SHAPE_TYPE.PICTURE):
      print("Found image")
      print("""Original file name """
            """{shape.image.filename}""")
      print("""Image type """
            f"""{shape.image.content_type}""")
      print("""Image extension """
            f"""{shape.image.ext}""")
      print("""Image size """
            f"""{shape.image.size}""")
      print("""Image dpi """
            f"""{shape.image.dpi}""")

      with open(shape.image.filename, 
                "wb") as f:
        f.write(shape.image.blob)                
            
# Find and extract all notes
print("""Slide notes""")
print("""===========""")
for slide in presentation.slides:
  if slide.has_notes_slide:
    s_n = presentation.slides.index(slide) + 1
    nts = (slide
           .notes_slide
           .notes_placeholder
           .text)
    print(f"""Notes for slide {s_n}:""")
    print(nts)