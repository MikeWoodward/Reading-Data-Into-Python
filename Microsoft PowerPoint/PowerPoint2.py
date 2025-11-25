#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Created on Mon Mar  3 15:05:47 2025

@author: mikewoodward
"""

# %%---------------------------------------------------------------------------
# Module metadata
# -----------------------------------------------------------------------------
__author__ = "mikewoodward"
__license__ = "MIT"
__summary__ = "One line summary"

# %%---------------------------------------------------------------------------
# Imports
# -----------------------------------------------------------------------------
from pptx import Presentation

if __name__ == "__main__":

    presentation = Presentation("Presentation2.pptx")

    replacement_text = {
        "#Company": "Cheapco",
        "#Date": "2025-09-12",
        "#Sales1": "Sales up from start of year",
        "#Sales2": "Q1 shows highest sales growth",
        "#Sales3": "Quarter end sales peak"}
    
    replacement_images = {
        "#CompanyLogo": "logo.png",
        "#SalesChart": "cheapco_sales_chart_quarterly.png"
        }

    # Iterate therough every slide
    for slide in presentation.slides:
        # Iterate through shapes in each slide
        for shape in slide.shapes:
            
            # Replace text
            if shape.has_text_frame:
                for old, new in replacement_text.items():
                    shape.text_frame.text = \
                        shape.text_frame.text.replace(old,
                                                      new)

            # Replace image
            alt_text = (shape
                        ._element
                        ._nvXxPr
                        .cNvPr
                        .attrib
                        .get("descr", ""))
            if alt_text in replacement_images:
                shape.element.getparent().remove(shape.element)
                slide.shapes.add_picture(
                    image_file=replacement_images[alt_text], 
                    left=shape.left, 
                    top=shape.top, 
                    width=shape.width, 
                    height=shape.height)

    presentation.save("Presentation3.pptx")
